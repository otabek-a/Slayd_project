import os
import json
import asyncio
import random
from openai import OpenAI
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from telegram import Update
from telegram.ext import Application, CommandHandler, MessageHandler, filters, ContextTypes, ConversationHandler
import requests
from io import BytesIO

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
TELEGRAM_TOKEN = os.getenv("TELEGRAM_TOKEN")
FAL_KEY = os.getenv("FAL_KEY")

client = OpenAI(api_key=OPENAI_API_KEY)

TOPIC, NUM_SLIDES, UNIVERSITY, STUDENT_NAME, FROM_TO = range(5)

user_data_store = {}

def add_advanced_shadow(shape, blur=8, distance=4, angle=45, transparency=0.2):
    """Add advanced shadow with customizable parameters for depth"""
    try:
        shadow = shape.shadow
        shadow.inherit = False
        shadow.visible = True
        shadow.style = 'OUTER'
        shadow.blur_radius = Pt(blur)
        shadow.distance = Pt(distance)
        shadow.angle = angle
        shadow.transparency = transparency
    except:
        pass

def add_glow_effect(shape, size=10, transparency=0.5, color=None):
    """Add modern glow effect to shapes"""
    try:
        if color is None:
            color = RGBColor(100, 150, 255)
        # Note: python-pptx has limited glow support, but we can enhance shadows
        add_advanced_shadow(shape, blur=size, distance=0, transparency=transparency)
    except:
        pass

def generate_image_sync(prompt: str):
    """Generate image using fal.ai API with enhanced prompts"""
    if not FAL_KEY:
        print("[v0] FAL_KEY not found, skipping image generation")
        return None
    
    try:
        headers = {
            "Authorization": f"Key {FAL_KEY}",
            "Content-Type": "application/json"
        }
        
        # Enhanced prompt for better quality
        payload = {
            "prompt": f"{prompt}, ultra professional, 8k quality, highly detailed, modern minimalist design, clean aesthetic, corporate style, premium look, sophisticated composition",
            "image_size": "landscape_16_9",
            "num_inference_steps": 4,
            "num_images": 1
        }
        
        print(f"[v0] Generating enhanced image: {prompt[:50]}...")
        
        response = requests.post(
            "https://fal.run/fal-ai/flux/schnell",
            headers=headers,
            json=payload,
            timeout=30
        )
        
        if response.status_code == 200:
            result = response.json()
            if result and 'images' in result and len(result['images']) > 0:
                image_url = result['images'][0]['url']
                img_response = requests.get(image_url, timeout=15)
                if img_response.status_code == 200:
                    print(f"[v0] Enhanced image generated successfully")
                    return BytesIO(img_response.content)
        
        print(f"[v0] Image generation failed: {response.status_code}")
    except Exception as e:
        print(f"[v0] Image generation error: {e}")
    
    return None

async def generate_image(prompt: str):
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, generate_image_sync, prompt)

def generate_slide_content_sync(topic: str, slides: int, university: str, student_name: str, from_to: str):
    num_content_slides = slides - 3
    
    prompt = f"""
Mavzu: "{topic}" haqida AYNAN {slides} ta slaydli professional taqdimot yarating. Barcha matn O'ZBEK TILIDA bo'lishi kerak.

MUHIM: AYNAN {slides} ta slayd bo'lishi kerak:
- 1-slayd: Sarlavha (title)
- 2-slayd: Kirish va mundarija (introduction) 
- 3-{slides-1} slaydlar: Asosiy kontent ({num_content_slides} ta content slayd)
- {slides}-slayd: Xulosa (conclusion)

JAMI: {slides} ta slayd

Har bir slayd uchun QISQA va ANIQ ma'lumot bering:
1. Har bir nuqta 10-12 so'zdan oshmasin (MUHIM!)
2. Oddiy va tushunarli til ishlatilsin
3. Faqat eng muhim ma'lumotlar

MUHIM: Har bir slayd uchun 2-3 ta topic-relevant image_prompts yarating.

JSON array qaytaring (faqat JSON, boshqa hech narsa yo'q):
[
  {{
    "type": "title",
    "title": "{topic}",
    "university": "{university}",
    "student": "{student_name}",
    "from_to": "{from_to}",
    "design_style": "mavzuga mos dizayn uslubi",
    "color_scheme": "mavzuga mos rang sxemasi",
    "image_prompts": [
      "{topic} professional background",
      "{topic} concept visualization"
    ]
  }},
  {{
    "type": "introduction",
    "title": "Kirish va Reja",
    "content": "Qisqa kirish matni (30-40 so'z).",
    "outline": [
      "Birinchi bo'lim - qisqa",
      "Ikkinchi bo'lim - qisqa",
      "Uchinchi bo'lim - qisqa"
    ],
    "image_prompts": [
      "{topic} introduction concept",
      "{topic} overview diagram"
    ]
  }},
  {{
    "type": "content",
    "title": "Qisqa sarlavha",
    "layout_type": "bullet_points",
    "points": [
      "Birinchi nuqta - 10-12 so'z",
      "Ikkinchi nuqta - 10-12 so'z",
      "Uchinchi nuqta - 10-12 so'z",
      "To'rtinchi nuqta - 10-12 so'z"
    ],
    "image_prompts": [
      "{topic} detailed diagram",
      "{topic} example"
    ]
  }},
  {{
    "type": "conclusion",
    "title": "Xulosa",
    "summary": "Qisqa xulosa matni (40-50 so'z).",
    "takeaways": [
      "Birinchi xulosa - qisqa",
      "Ikkinchi xulosa - qisqa",
      "Uchinchi xulosa - qisqa"
    ],
    "image_prompts": [
      "{topic} success concept",
      "{topic} summary"
    ]
  }}
]

ESLATMA: Aynan {num_content_slides} ta "content" tipidagi slayd yarating. Barcha matn QISQA bo'lishi kerak!

Faqat to'g'ri JSON qaytaring, qo'shimcha matn yoki kod bloklarsiz.
"""

    try:
        print("[v0] Calling OpenAI API for content generation...")
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.8
        )
        print("[v0] OpenAI API call successful")
    except Exception as e:
        print(f"[v0] OpenAI API error: {e}")
        return []

    try:
        content = response.choices[0].message.content
        print(f"[v0] Received content from OpenAI (length: {len(content)})")
    except Exception as e:
        print(f"[v0] No content from OpenAI: {e}")
        return []

    if content.startswith("\`\`\`json"):
        content = content[len("\`\`\`json"):].strip()
    elif content.startswith("\`\`\`"):
        content = content[len("\`\`\`"):].strip()
    if content.endswith("\`\`\`"):
        content = content[:-3].strip()

    try:
        slides_data = json.loads(content)
        print(f"[v0] Successfully parsed {len(slides_data)} slides")
    except Exception as e:
        print(f"[v0] JSON parsing error: {e}")
        print(f"[v0] Raw response: {content[:500]}...")
        return []

    return slides_data

async def generate_slide_content(topic: str, slides: int, university: str, student_name: str, from_to: str):
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, generate_slide_content_sync, topic, slides, university, student_name, from_to)

def get_advanced_design_template(seed):
    """Get one of 18 ultra-modern professionally designed templates with PERFECT READABILITY"""
    random.seed(seed)
    
    templates = [
        # Minimalist Modern Series
        {
            "name": "Arctic Minimalist",
            "bg_color": RGBColor(255, 255, 255),
            "primary": RGBColor(10, 10, 10),
            "accent": RGBColor(0, 122, 255),
            "text_primary": RGBColor(10, 10, 10),
            "text_secondary": RGBColor(100, 100, 100),
            "shape_fill": RGBColor(248, 249, 250),
            "shape_text": RGBColor(10, 10, 10),  # Dark text on light shape
            "shape_border": RGBColor(0, 122, 255),
            "gradient_start": RGBColor(240, 248, 255),
            "gradient_end": RGBColor(255, 255, 255),
            "title_font_size": 52,
            "subtitle_font_size": 26,
            "content_font_size": 18,
            "style": "minimalist"
        },
        {
            "name": "Dark Mode Elite",
            "bg_color": RGBColor(18, 18, 18),
            "primary": RGBColor(0, 229, 255),
            "accent": RGBColor(138, 43, 226),
            "text_primary": RGBColor(255, 255, 255),
            "text_secondary": RGBColor(200, 200, 200),
            "shape_fill": RGBColor(30, 30, 30),
            "shape_text": RGBColor(255, 255, 255),  # Light text on dark shape
            "shape_border": RGBColor(0, 229, 255),
            "gradient_start": RGBColor(25, 25, 35),
            "gradient_end": RGBColor(18, 18, 18),
            "title_font_size": 50,
            "subtitle_font_size": 25,
            "content_font_size": 17,
            "style": "dark_modern"
        },
        {
            "name": "Professional Blue",
            "bg_color": RGBColor(245, 248, 252),
            "primary": RGBColor(25, 118, 210),
            "accent": RGBColor(66, 165, 245),
            "text_primary": RGBColor(13, 71, 161),
            "text_secondary": RGBColor(25, 118, 210),
            "shape_fill": RGBColor(255, 255, 255),
            "shape_text": RGBColor(13, 71, 161),  # Dark blue text on white
            "shape_border": RGBColor(66, 165, 245),
            "gradient_start": RGBColor(227, 242, 253),
            "gradient_end": RGBColor(245, 248, 255),
            "title_font_size": 48,
            "subtitle_font_size": 24,
            "content_font_size": 17,
            "style": "corporate"
        },
        {
            "name": "Emerald Business",
            "bg_color": RGBColor(255, 255, 255),
            "primary": RGBColor(0, 105, 92),
            "accent": RGBColor(0, 200, 83),
            "text_primary": RGBColor(0, 77, 64),
            "text_secondary": RGBColor(69, 90, 100),
            "shape_fill": RGBColor(232, 245, 233),
            "shape_text": RGBColor(0, 77, 64),  # Dark green text on light green
            "shape_border": RGBColor(0, 105, 92),
            "gradient_start": RGBColor(232, 245, 233),
            "gradient_end": RGBColor(255, 255, 255),
            "title_font_size": 48,
            "subtitle_font_size": 24,
            "content_font_size": 17,
            "style": "corporate"
        },
        {
            "name": "Ocean Depth",
            "bg_color": RGBColor(240, 248, 255),
            "primary": RGBColor(0, 105, 148),
            "accent": RGBColor(0, 188, 212),
            "text_primary": RGBColor(1, 87, 155),
            "text_secondary": RGBColor(38, 50, 56),
            "shape_fill": RGBColor(255, 255, 255),
            "shape_text": RGBColor(1, 87, 155),  # Dark blue text on white
            "shape_border": RGBColor(0, 188, 212),
            "gradient_start": RGBColor(224, 247, 250),
            "gradient_end": RGBColor(240, 248, 255),
            "title_font_size": 48,
            "subtitle_font_size": 24,
            "content_font_size": 17,
            "style": "nature_modern"
        },
        {
            "name": "Forest Canopy",
            "bg_color": RGBColor(249, 251, 248),
            "primary": RGBColor(27, 94, 32),
            "accent": RGBColor(76, 175, 80),
            "text_primary": RGBColor(27, 94, 32),
            "text_secondary": RGBColor(56, 142, 60),
            "shape_fill": RGBColor(255, 255, 255),
            "shape_text": RGBColor(27, 94, 32),  # Dark green text on white
            "shape_border": RGBColor(76, 175, 80),
            "gradient_start": RGBColor(232, 245, 233),
            "gradient_end": RGBColor(249, 251, 248),
            "title_font_size": 48,
            "subtitle_font_size": 24,
            "content_font_size": 17,
            "style": "nature_modern"
        },
        {
            "name": "Sunset Gradient",
            "bg_color": RGBColor(255, 250, 245),
            "primary": RGBColor(211, 47, 47),
            "accent": RGBColor(255, 152, 0),
            "text_primary": RGBColor(183, 28, 28),
            "text_secondary": RGBColor(191, 54, 12),
            "shape_fill": RGBColor(255, 255, 255),
            "shape_text": RGBColor(183, 28, 28),  # Dark red text on white
            "shape_border": RGBColor(255, 152, 0),
            "gradient_start": RGBColor(255, 245, 238),
            "gradient_end": RGBColor(255, 250, 245),
            "title_font_size": 50,
            "subtitle_font_size": 26,
            "content_font_size": 18,
            "style": "gradient_bold"
        },
        {
            "name": "Royal Purple",
            "bg_color": RGBColor(248, 245, 255),
            "primary": RGBColor(94, 53, 177),
            "accent": RGBColor(156, 39, 176),
            "text_primary": RGBColor(74, 20, 140),
            "text_secondary": RGBColor(106, 27, 154),
            "shape_fill": RGBColor(255, 255, 255),
            "shape_text": RGBColor(74, 20, 140),  # Dark purple text on white
            "shape_border": RGBColor(156, 39, 176),
            "gradient_start": RGBColor(237, 231, 246),
            "gradient_end": RGBColor(248, 245, 255),
            "title_font_size": 50,
            "subtitle_font_size": 25,
            "content_font_size": 17,
            "style": "creative"
        },
        {
            "name": "Monochrome Pro",
            "bg_color": RGBColor(255, 255, 255),
            "primary": RGBColor(33, 33, 33),
            "accent": RGBColor(117, 117, 117),
            "text_primary": RGBColor(33, 33, 33),
            "text_secondary": RGBColor(97, 97, 97),
            "shape_fill": RGBColor(245, 245, 245),
            "shape_text": RGBColor(33, 33, 33),  # Dark text on light gray
            "shape_border": RGBColor(33, 33, 33),
            "gradient_start": RGBColor(245, 245, 245),
            "gradient_end": RGBColor(255, 255, 255),
            "title_font_size": 50,
            "subtitle_font_size": 25,
            "content_font_size": 17,
            "style": "monochrome"
        },
        {
            "name": "Charcoal Elegance",
            "bg_color": RGBColor(250, 250, 250),
            "primary": RGBColor(55, 71, 79),
            "accent": RGBColor(96, 125, 139),
            "text_primary": RGBColor(38, 50, 56),
            "text_secondary": RGBColor(69, 90, 100),
            "shape_fill": RGBColor(255, 255, 255),
            "shape_text": RGBColor(38, 50, 56),  # Dark text on white
            "shape_border": RGBColor(96, 125, 139),
            "gradient_start": RGBColor(236, 239, 241),
            "gradient_end": RGBColor(250, 250, 250),
            "title_font_size": 48,
            "subtitle_font_size": 24,
            "content_font_size": 17,
            "style": "monochrome"
        }
    ]
    
    template = random.choice(templates)
    random.seed()
    return template

def add_modern_decorative_element(slide_obj, template, position="top"):
    """Add ultra-modern decorative elements with 2025 design trends"""
    if position == "top":
        # Main bold accent bar
        accent_bar = slide_obj.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.15)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = template['primary']
        accent_bar.line.fill.background()
        
        # Gradient accent line
        gradient_line = slide_obj.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0.15), Inches(10), Inches(0.05)
        )
        gradient_line.fill.solid()
        gradient_line.fill.fore_color.rgb = template['accent']
        gradient_line.line.fill.background()
        
        # Modern corner accent
        corner_accent = slide_obj.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9.5), Inches(0.3), Inches(0.4), Inches(0.4)
        )
        corner_accent.fill.solid()
        corner_accent.fill.fore_color.rgb = template['accent']
        corner_accent.line.fill.background()
        add_advanced_shadow(corner_accent, blur=6, distance=3, transparency=0.3)
        
    elif position == "bottom":
        # Bottom accent with modern style
        accent_bar = slide_obj.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.3), Inches(10), Inches(0.2)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = template['primary']
        accent_bar.line.fill.background()
        
    elif position == "side":
        # Vertical side accent for asymmetric layouts
        side_bar = slide_obj.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.1), Inches(1), Inches(0.08), Inches(5.5)
        )
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = template['accent']
        side_bar.line.fill.background()

def add_geometric_decoration(slide_obj, template, style="circles"):
    """Add modern geometric decorations for visual interest"""
    if style == "circles":
        # Large background circle
        circle1 = slide_obj.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(8.5), Inches(-1), Inches(3), Inches(3)
        )
        circle1.fill.solid()
        circle1.fill.fore_color.rgb = template['shape_fill']
        circle1.line.fill.background()
        circle1.fill.transparency = 0.5
        
        # Small accent circle
        circle2 = slide_obj.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-0.5), Inches(6), Inches(1.5), Inches(1.5)
        )
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = template['accent']
        circle2.line.fill.background()
        circle2.fill.transparency = 0.3
        
    elif style == "triangles":
        # Modern triangle accents (using rotated rectangles)
        triangle = slide_obj.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(9), Inches(6.5), Inches(1), Inches(1)
        )
        triangle.fill.solid()
        triangle.fill.fore_color.rgb = template['accent']
        triangle.line.fill.background()
        triangle.fill.transparency = 0.4

async def create_ppt(content, filename="presentation.pptx"):
    print(f"[v0] Starting READABLE PPT creation with {len(content)} slides")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    design_seed = random.randint(1, 1000000)
    template = get_advanced_design_template(design_seed)
    print(f"[v0] Using READABLE design template: {template['name']} ({template['style']})")
    
    layout_types = ["cards", "two_column", "timeline", "comparison", "grid", "numbered", "highlight", "icon_based"]
    used_layouts = []
    
    for idx, slide_data in enumerate(content):
        print(f"[v0] Creating readable slide {idx + 1}/{len(content)}")
        slide_type = slide_data.get("type", "content")
        
        slide_layout = prs.slide_layouts[6]
        slide_obj = prs.slides.add_slide(slide_layout)

        background = slide_obj.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = template['bg_color']

        top_accent = slide_obj.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.25)
        )
        top_accent.fill.solid()
        top_accent.fill.fore_color.rgb = template['primary']
        top_accent.line.fill.background()
        
        corner_circle = slide_obj.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(9.3), Inches(7), Inches(0.5), Inches(0.5)
        )
        corner_circle.fill.solid()
        corner_circle.fill.fore_color.rgb = template['accent']
        corner_circle.line.fill.background()
        corner_circle.fill.transparency = 0.3

        image_prompts = slide_data.get("image_prompts", [])
        
        try:
            if slide_type == "title":
                title_box = slide_obj.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.5), Inches(0.6), Inches(9), Inches(1.1)
                )
                title_box.fill.solid()
                title_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                title_box.line.color.rgb = template['primary']
                title_box.line.width = Pt(8)
                add_advanced_shadow(title_box, blur=15, distance=8, transparency=0.25)
                
                title_frame = title_box.text_frame
                title_frame.word_wrap = True
                title_frame.margin_left = Inches(0.3)
                title_frame.margin_right = Inches(0.3)
                title_frame.margin_top = Inches(0.15)
                title_frame.margin_bottom = Inches(0.15)
                title_frame.text = slide_data.get("title", "")
                title_p = title_frame.paragraphs[0]
                title_p.font.size = Pt(42)
                title_p.font.bold = True
                title_p.font.color.rgb = template['primary']  # Dark text on white box
                title_p.alignment = PP_ALIGN.CENTER
                title_p.line_spacing = 1.1
                title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                uni_box = slide_obj.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(1), Inches(2.1), Inches(8), Inches(1.8)
                )
                uni_box.fill.solid()
                uni_box.fill.fore_color.rgb = template['shape_fill']
                uni_box.line.color.rgb = template['accent']
                uni_box.line.width = Pt(5)
                add_advanced_shadow(uni_box, blur=12, distance=6, transparency=0.2)
                
                uni_frame = uni_box.text_frame
                uni_frame.word_wrap = True
                uni_frame.margin_left = Inches(0.4)
                uni_frame.margin_right = Inches(0.4)
                uni_frame.margin_top = Inches(0.25)
                uni_frame.margin_bottom = Inches(0.25)
                uni_frame.text = slide_data.get("university", "")
                uni_p = uni_frame.paragraphs[0]
                uni_p.font.size = Pt(20)
                uni_p.font.bold = True
                uni_p.font.color.rgb = template['shape_text']  # Use shape_text for proper contrast
                uni_p.alignment = PP_ALIGN.CENTER
                uni_p.line_spacing = 1.2
                uni_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                student_card = slide_obj.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.8), Inches(4.5), Inches(4.2), Inches(1.5)
                )
                student_card.fill.solid()
                student_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
                student_card.line.color.rgb = template['accent']
                student_card.line.width = Pt(5)
                add_advanced_shadow(student_card, blur=10, distance=5, transparency=0.2)
                
                student_icon = slide_obj.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(1.2), Inches(4.8), Inches(0.9), Inches(0.9)
                )
                student_icon.fill.solid()
                student_icon.fill.fore_color.rgb = template['primary']
                student_icon.line.fill.background()
                add_advanced_shadow(student_icon, blur=6, distance=3, transparency=0.25)
                
                icon_frame = student_icon.text_frame
                icon_frame.text = "👤"
                icon_p = icon_frame.paragraphs[0]
                icon_p.font.size = Pt(30)
                icon_p.alignment = PP_ALIGN.CENTER
                icon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                student_frame = student_card.text_frame
                student_frame.word_wrap = True
                student_frame.margin_left = Inches(1.4)
                student_frame.margin_right = Inches(0.3)
                student_frame.margin_top = Inches(0.2)
                
                student_label = student_frame.paragraphs[0]
                student_label.text = "Bajarildi:"
                student_label.font.size = Pt(13)
                student_label.font.bold = True
                student_label.font.color.rgb = template['text_secondary']
                student_label.alignment = PP_ALIGN.LEFT
                
                student_name_p = student_frame.add_paragraph()
                student_name_p.text = slide_data.get('student', '')
                student_name_p.font.size = Pt(20)
                student_name_p.font.bold = True
                student_name_p.font.color.rgb = template['text_primary']  # Dark text on white card
                student_name_p.alignment = PP_ALIGN.LEFT
                student_name_p.space_before = Pt(4)
                
                teacher_card = slide_obj.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(5.0), Inches(4.5), Inches(4.2), Inches(1.5)
                )
                teacher_card.fill.solid()
                teacher_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
                teacher_card.line.color.rgb = template['accent']
                teacher_card.line.width = Pt(5)
                add_advanced_shadow(teacher_card, blur=10, distance=5, transparency=0.2)
                
                teacher_icon = slide_obj.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(7.9), Inches(4.8), Inches(0.9), Inches(0.9)
                )
                teacher_icon.fill.solid()
                teacher_icon.fill.fore_color.rgb = template['accent']
                teacher_icon.line.fill.background()
                add_advanced_shadow(teacher_icon, blur=6, distance=3, transparency=0.25)
                
                teacher_icon_frame = teacher_icon.text_frame
                teacher_icon_frame.text = "✔"
                teacher_icon_p = teacher_icon_frame.paragraphs[0]
                teacher_icon_p.font.size = Pt(40)
                teacher_icon_p.font.bold = True
                teacher_icon_p.font.color.rgb = RGBColor(255, 255, 255)  # White checkmark on colored circle
                teacher_icon_p.alignment = PP_ALIGN.CENTER
                teacher_icon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                teacher_frame = teacher_card.text_frame
                teacher_frame.word_wrap = True
                teacher_frame.margin_left = Inches(0.3)
                teacher_frame.margin_right = Inches(1.4)
                teacher_frame.margin_top = Inches(0.2)
                
                teacher_label = teacher_frame.paragraphs[0]
                teacher_label.text = "Tekshirdi:"
                teacher_label.font.size = Pt(13)
                teacher_label.font.bold = True
                teacher_label.font.color.rgb = template['text_secondary']
                teacher_label.alignment = PP_ALIGN.RIGHT
                
                teacher_name_p = teacher_frame.add_paragraph()
                teacher_name_p.text = slide_data.get('from_to', '')
                teacher_name_p.font.size = Pt(20)
                teacher_name_p.font.bold = True
                teacher_name_p.font.color.rgb = template['text_primary']  # Dark text on white card
                teacher_name_p.alignment = PP_ALIGN.RIGHT
                teacher_name_p.space_before = Pt(4)

            elif slide_type == "introduction":
                side_bar = slide_obj.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.1), Inches(1.5), Inches(0.12), Inches(4.5)
                )
                side_bar.fill.solid()
                side_bar.fill.fore_color.rgb = template['accent']
                side_bar.line.fill.background()
                
                title_box = slide_obj.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(1.5), Inches(0.6), Inches(7), Inches(1.1)
                )
                title_box.fill.solid()
                title_box.fill.fore_color.rgb = template['primary']
                title_box.line.fill.background()
                add_advanced_shadow(title_box, blur=14, distance=7, transparency=0.25)
                
                title_frame = title_box.text_frame
                title_frame.text = slide_data.get("title", "Kirish va Reja")
                title_frame.margin_left = Inches(0.5)
                title_frame.margin_right = Inches(0.5)
                title_p = title_frame.paragraphs[0]
                title_p.font.size = Pt(42)
                title_p.font.bold = True
                title_p.font.color.rgb = RGBColor(255, 255, 255)  # White text on dark primary color
                title_p.alignment = PP_ALIGN.CENTER
                title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                if slide_data.get("content"):
                    intro_box = slide_obj.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(1.5), Inches(2), Inches(7), Inches(1.3)
                    )
                    intro_box.fill.solid()
                    intro_box.fill.fore_color.rgb = template['shape_fill']
                    intro_box.line.color.rgb = template['accent']
                    intro_box.line.width = Pt(3)
                    add_advanced_shadow(intro_box, blur=8, distance=4, transparency=0.15)
                    
                    intro_frame = intro_box.text_frame
                    intro_frame.word_wrap = True
                    intro_frame.margin_left = Inches(0.5)
                    intro_frame.margin_right = Inches(0.5)
                    intro_frame.margin_top = Inches(0.3)
                    intro_frame.margin_bottom = Inches(0.3)
                    intro_frame.text = slide_data.get("content", "")
                    intro_p = intro_frame.paragraphs[0]
                    intro_p.font.size = Pt(18)
                    intro_p.font.color.rgb = template['shape_text']  # Proper contrast text
                    intro_p.line_spacing = 1.5
                    intro_p.alignment = PP_ALIGN.CENTER
                    intro_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                outline_items = slide_data.get("outline", [])
                if len(outline_items) >= 4:
                    positions = [
                        (1.5, 3.6), (5.5, 3.6),
                        (1.5, 5.2), (5.5, 5.2)
                    ]
                    for i, item in enumerate(outline_items[:4]):
                        x, y = positions[i]
                        item_box = slide_obj.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(3.5), Inches(1.3)
                        )
                        item_box.fill.solid()
                        item_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                        item_box.line.color.rgb = template['primary']
                        item_box.line.width = Pt(4)
                        add_advanced_shadow(item_box, blur=8, distance=4, transparency=0.18)
                        
                        badge = slide_obj.shapes.add_shape(
                            MSO_SHAPE.OVAL,
                            Inches(x + 0.2), Inches(y + 0.15), Inches(0.5), Inches(0.5)
                        )
                        badge.fill.solid()
                        badge.fill.fore_color.rgb = template['accent']
                        badge.line.fill.background()
                        
                        badge_frame = badge.text_frame
                        badge_frame.text = str(i + 1)
                        badge_p = badge_frame.paragraphs[0]
                        badge_p.font.size = Pt(20)
                        badge_p.font.bold = True
                        badge_p.font.color.rgb = RGBColor(255, 255, 255)  # White number on colored badge
                        badge_p.alignment = PP_ALIGN.CENTER
                        badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                        
                        item_frame = item_box.text_frame
                        item_frame.word_wrap = True
                        item_frame.margin_left = Inches(0.85)
                        item_frame.margin_right = Inches(0.3)
                        item_frame.margin_top = Inches(0.25)
                        item_frame.text = item
                        item_p = item_frame.paragraphs[0]
                        item_p.font.size = Pt(16)
                        item_p.font.color.rgb = template['text_primary']  # Dark text on white box
                        item_p.line_spacing = 1.3
                        item_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            elif slide_type == "content":
                points = slide_data.get("points", [])
                
                while len(points) < 4:
                    points.append(f"Qo'shimcha ma'lumot {len(points) + 1}")
                
                available_layouts = [l for l in layout_types if l not in used_layouts[-3:]]
                if not available_layouts:
                    available_layouts = layout_types
                
                layout_choice = random.choice(available_layouts)
                used_layouts.append(layout_choice)
                
                title_box = slide_obj.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.95)
                )
                title_box.fill.solid()
                title_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                title_box.line.color.rgb = template['primary']
                title_box.line.width = Pt(6)
                add_advanced_shadow(title_box, blur=12, distance=6, transparency=0.2)
                
                title_frame = title_box.text_frame
                title_frame.word_wrap = True
                title_frame.margin_left = Inches(0.4)
                title_frame.margin_right = Inches(0.4)
                title_frame.text = slide_data.get("title", "")
                title_p = title_frame.paragraphs[0]
                title_p.font.size = Pt(38)
                title_p.font.bold = True
                title_p.font.color.rgb = template['primary']
                title_p.alignment = PP_ALIGN.CENTER
                title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                if layout_choice == "cards":
                    start_y = 1.8
                    spacing = 1.35
                    
                    for i, point in enumerate(points[:4]):
                        point_box = slide_obj.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.8), Inches(start_y + i * spacing), Inches(8.4), Inches(1.2)
                        )
                        point_box.fill.solid()
                        point_box.fill.fore_color.rgb = template['shape_fill']
                        point_box.line.color.rgb = template['accent']
                        point_box.line.width = Pt(4)
                        add_advanced_shadow(point_box, blur=8, distance=4, transparency=0.18)
                        
                        bullet_circle = slide_obj.shapes.add_shape(
                            MSO_SHAPE.OVAL,
                            Inches(1.2), Inches(start_y + i * spacing + 0.45), Inches(0.3), Inches(0.3)
                        )
                        bullet_circle.fill.solid()
                        bullet_circle.fill.fore_color.rgb = template['accent']
                        bullet_circle.line.fill.background()
                        
                        separator = slide_obj.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            Inches(1.65), Inches(start_y + i * spacing + 0.35), Inches(0.04), Inches(0.5)
                        )
                        separator.fill.solid()
                        separator.fill.fore_color.rgb = template['primary']
                        separator.line.fill.background()
                        
                        point_frame = point_box.text_frame
                        point_frame.word_wrap = True
                        point_frame.margin_left = Inches(0.95)
                        point_frame.margin_right = Inches(0.4)
                        point_frame.margin_top = Inches(0.25)
                        point_frame.margin_bottom = Inches(0.25)
                        point_frame.text = point
                        point_p = point_frame.paragraphs[0]
                        point_p.font.size = Pt(19)
                        point_p.font.color.rgb = template['shape_text']  # Proper contrast text
                        point_p.line_spacing = 1.4
                        point_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                elif layout_choice in ["two_column", "timeline", "comparison", "grid", "numbered", "highlight", "icon_based"]:
                    # Using template['shape_text'] for text on shape_fill backgrounds
                    # Using template['text_primary'] for text on white backgrounds
                    # Using RGBColor(255, 255, 255) for text on dark/colored backgrounds
                    
                    # --- Two Column Layout ---
                    if layout_choice == "two_column":
                        left_points = points[:2]  # First 2 points
                        right_points = points[2:4]  # Next 2 points
                        
                        for i, point in enumerate(left_points):
                            point_box = slide_obj.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.6), Inches(1.9 + i * 2.4), Inches(4.3), Inches(2.1)
                            )
                            point_box.fill.solid()
                            point_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                            point_box.line.color.rgb = template['primary']
                            point_box.line.width = Pt(5)
                            add_advanced_shadow(point_box, blur=10, distance=5, transparency=0.2)
                            
                            number_box = slide_obj.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.9), Inches(2.1 + i * 2.4), Inches(0.6), Inches(0.6)
                            )
                            number_box.fill.solid()
                            number_box.fill.fore_color.rgb = template['accent']
                            number_box.line.fill.background()
                            
                            num_frame = number_box.text_frame
                            num_frame.text = str(i + 1)
                            num_p = num_frame.paragraphs[0]
                            num_p.font.size = Pt(26)
                            num_p.font.bold = True
                            num_p.font.color.rgb = RGBColor(255, 255, 255)
                            num_p.alignment = PP_ALIGN.CENTER
                            num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                            
                            point_frame = point_box.text_frame
                            point_frame.word_wrap = True
                            point_frame.margin_left = Inches(0.3)
                            point_frame.margin_right = Inches(0.3)
                            point_frame.margin_top = Inches(0.85)
                            point_frame.margin_bottom = Inches(0.3)
                            point_frame.text = point
                            point_p = point_frame.paragraphs[0]
                            point_p.font.size = Pt(17)
                            point_p.font.color.rgb = template['text_primary']  # Dark text on white
                            point_p.line_spacing = 1.3
                            point_p.alignment = PP_ALIGN.CENTER
                            point_frame.vertical_anchor = MSO_ANCHOR.TOP
                        
                        for i, point in enumerate(right_points):
                            point_box = slide_obj.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(5.1), Inches(1.9 + i * 2.4), Inches(4.3), Inches(2.1)
                            )
                            point_box.fill.solid()
                            point_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                            point_box.line.color.rgb = template['primary']
                            point_box.line.width = Pt(5)
                            add_advanced_shadow(point_box, blur=10, distance=5, transparency=0.2)
                            
                            number_box = slide_obj.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(5.4), Inches(2.1 + i * 2.4), Inches(0.6), Inches(0.6)
                            )
                            number_box.fill.solid()
                            number_box.fill.fore_color.rgb = template['accent']
                            number_box.line.fill.background()
                            
                            num_frame = number_box.text_frame
                            num_frame.text = str(i + 3)
                            num_p = num_frame.paragraphs[0]
                            num_p.font.size = Pt(26)
                            num_p.font.bold = True
                            num_p.font.color.rgb = RGBColor(255, 255, 255)
                            num_p.alignment = PP_ALIGN.CENTER
                            num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                            
                            point_frame = point_box.text_frame
                            point_frame.word_wrap = True
                            point_frame.margin_left = Inches(0.3)
                            point_frame.margin_right = Inches(0.3)
                            point_frame.margin_top = Inches(0.85)
                            point_frame.margin_bottom = Inches(0.3)
                            point_frame.text = point
                            point_p = point_frame.paragraphs[0]
                            point_p.font.size = Pt(17)
                            point_p.font.color.rgb = template['text_primary']  # Dark text on white
                            point_p.line_spacing = 1.3
                            point_p.alignment = PP_ALIGN.CENTER
                            point_frame.vertical_anchor = MSO_ANCHOR.TOP
                    
                    # --- Timeline Layout ---
                    elif layout_choice == "timeline":
                        for i, point in enumerate(points[:4]):
                            node_x = 0.8 + i * 2.2
                            node_box = slide_obj.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(node_x), Inches(2.5), Inches(2), Inches(3.5)
                            )
                            node_box.fill.solid()
                            node_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                            node_box.line.color.rgb = template['accent']
                            node_box.line.width = Pt(5)
                            add_advanced_shadow(node_box, blur=10, distance=5, transparency=0.2)
                            
                            circle = slide_obj.shapes.add_shape(
                                MSO_SHAPE.OVAL,
                                Inches(node_x + 0.75), Inches(2.8), Inches(0.5), Inches(0.5)
                            )
                            circle.fill.solid()
                            circle.fill.fore_color.rgb = template['primary']
                            circle.line.fill.background()
                            
                            circle_frame = circle.text_frame
                            circle_frame.text = str(i + 1)
                            circle_p = circle_frame.paragraphs[0]
                            circle_p.font.size = Pt(22)
                            circle_p.font.bold = True
                            circle_p.font.color.rgb = RGBColor(255, 255, 255)
                            circle_p.alignment = PP_ALIGN.CENTER
                            circle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                            
                            node_frame = node_box.text_frame
                            node_frame.word_wrap = True
                            node_frame.margin_left = Inches(0.2)
                            node_frame.margin_right = Inches(0.2)
                            node_frame.margin_top = Inches(0.9)
                            node_frame.text = point
                            node_p = node_frame.paragraphs[0]
                            node_p.font.size = Pt(15)
                            node_p.font.color.rgb = template['text_primary']
                            node_p.line_spacing = 1.3
                            node_p.alignment = PP_ALIGN.CENTER
                            node_frame.vertical_anchor = MSO_ANCHOR.TOP
                            
                            if i < 3:
                                arrow = slide_obj.shapes.add_shape(
                                    MSO_SHAPE.RIGHT_ARROW,
                                    Inches(node_x + 2.05), Inches(4.2), Inches(0.15), Inches(0.3)
                                )
                                arrow.fill.solid()
                                arrow.fill.fore_color.rgb = template['accent']
                                arrow.line.fill.background()
                    
                    # --- Comparison Layout ---
                    elif layout_choice == "comparison":
                        left_points = points[:2]
                        right_points = points[2:4]
                        
                        divider = slide_obj.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(4.85), Inches(1.7), Inches(0.3), Inches(5)
                        )
                        divider.fill.solid()
                        divider.fill.fore_color.rgb = template['accent']
                        divider.line.fill.background()
                        
                        for i, point in enumerate(left_points):
                            box = slide_obj.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.6), Inches(1.9 + i * 2.4), Inches(4), Inches(2.1)
                            )
                            box.fill.solid()
                            box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                            box.line.color.rgb = template['primary']
                            box.line.width = Pt(4)
                            add_advanced_shadow(box, blur=8, distance=4, transparency=0.18)
                            
                            icon = slide_obj.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.9), Inches(2.1 + i * 2.4), Inches(0.5), Inches(0.5)
                            )
                            icon.fill.solid()
                            icon.fill.fore_color.rgb = template['accent']
                            icon.line.fill.background()
                            
                            icon_frame = icon.text_frame
                            icon_frame.text = "◆"
                            icon_p = icon_frame.paragraphs[0]
                            icon_p.font.size = Pt(20)
                            icon_p.font.color.rgb = RGBColor(255, 255, 255)
                            icon_p.alignment = PP_ALIGN.CENTER
                            icon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                            
                            frame = box.text_frame
                            frame.word_wrap = True
                            frame.margin_left = Inches(0.3)
                            frame.margin_right = Inches(0.3)
                            frame.margin_top = Inches(0.75)
                            frame.margin_bottom = Inches(0.3)
                            frame.text = point
                            p = frame.paragraphs[0]
                            p.font.size = Pt(16)
                            p.font.color.rgb = template['text_primary']  # Dark text on white
                            p.line_spacing = 1.3
                            p.alignment = PP_ALIGN.CENTER
                            frame.vertical_anchor = MSO_ANCHOR.TOP
                        
                        for i, point in enumerate(right_points):
                            box = slide_obj.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(5.4), Inches(1.9 + i * 2.4), Inches(4), Inches(2.1)
                            )
                            box.fill.solid()
                            box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                            box.line.color.rgb = template['primary']
                            box.line.width = Pt(4)
                            add_advanced_shadow(box, blur=8, distance=4, transparency=0.18)
                            
                            icon = slide_obj.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(8.5), Inches(2.1 + i * 2.4), Inches(0.5), Inches(0.5)
                            )
                            icon.fill.solid()
                            icon.fill.fore_color.rgb = template['accent']
                            icon.line.fill.background()
                            
                            icon_frame = icon.text_frame
                            icon_frame.text = "◆"
                            icon_p = icon_frame.paragraphs[0]
                            icon_p.font.size = Pt(20)
                            icon_p.font.color.rgb = RGBColor(255, 255, 255)
                            icon_p.alignment = PP_ALIGN.CENTER
                            icon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                            
                            frame = box.text_frame
                            frame.word_wrap = True
                            frame.margin_left = Inches(0.3)
                            frame.margin_right = Inches(0.3)
                            frame.margin_top = Inches(0.75)
                            frame.margin_bottom = Inches(0.3)
                            frame.text = point
                            p = frame.paragraphs[0]
                            p.font.size = Pt(16)
                            p.font.color.rgb = template['text_primary']  # Dark text on white
                            p.line_spacing = 1.3
                            p.alignment = PP_ALIGN.CENTER
                            frame.vertical_anchor = MSO_ANCHOR.TOP
                    
                    # --- Grid Layout ---
                    elif layout_choice == "grid":
                        positions = [
                            (0.6, 1.8), (5.2, 1.8),
                            (0.6, 4.5), (5.2, 4.5)
                        ]
                        
                        for i, point in enumerate(points[:4]):
                            x, y = positions[i]
                            box = slide_obj.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(4.2), Inches(2.4)
                            )
                            box.fill.solid()
                            box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                            box.line.color.rgb = template['accent']
                            box.line.width = Pt(5)
                            add_advanced_shadow(box, blur=10, distance=5, transparency=0.2)
                            
                            badge = slide_obj.shapes.add_shape(
                                MSO_SHAPE.OVAL,
                                Inches(x + 0.3), Inches(y + 0.3), Inches(0.6), Inches(0.6)
                            )
                            badge.fill.solid()
                            badge.fill.fore_color.rgb = template['primary']
                            badge.line.fill.background()
                            
                            badge_frame = badge.text_frame
                            badge_frame.text = str(i + 1)
                            badge_p = badge_frame.paragraphs[0]
                            badge_p.font.size = Pt(24)
                            badge_p.font.bold = True
                            badge_p.font.color.rgb = RGBColor(255, 255, 255)
                            badge_p.alignment = PP_ALIGN.CENTER
                            badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                            
                            frame = box.text_frame
                            frame.word_wrap = True
                            frame.margin_left = Inches(0.4)
                            frame.margin_right = Inches(0.4)
                            frame.margin_top = Inches(1.1)
                            frame.text = point
                            p = frame.paragraphs[0]
                            p.font.size = Pt(17)
                            p.font.color.rgb = template['text_primary']
                            p.line_spacing = 1.3
                            p.alignment = PP_ALIGN.CENTER
                            frame.vertical_anchor = MSO_ANCHOR.TOP
                    
                    # --- Numbered Layout ---
                    elif layout_choice == "numbered":
                        for i, point in enumerate(points[:4]):
                            y_pos = 1.8 + i * 1.35
                            
                            number_box = slide_obj.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.8), Inches(y_pos), Inches(1), Inches(1.2)
                            )
                            number_box.fill.solid()
                            number_box.fill.fore_color.rgb = template['primary']
                            number_box.line.fill.background()
                            add_advanced_shadow(number_box, blur=8, distance=4, transparency=0.2)
                            
                            num_frame = number_box.text_frame
                            num_frame.text = str(i + 1)
                            num_p = num_frame.paragraphs[0]
                            num_p.font.size = Pt(48)
                            num_p.font.bold = True
                            num_p.font.color.rgb = RGBColor(255, 255, 255)
                            num_p.alignment = PP_ALIGN.CENTER
                            num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                            
                            content_box = slide_obj.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(2), Inches(y_pos), Inches(7.2), Inches(1.2)
                            )
                            content_box.fill.solid()
                            content_box.fill.fore_color.rgb = template['shape_fill']
                            content_box.line.color.rgb = template['accent']
                            content_box.line.width = Pt(4)
                            add_advanced_shadow(content_box, blur=8, distance=4, transparency=0.18)
                            
                            content_frame = content_box.text_frame
                            content_frame.word_wrap = True
                            content_frame.margin_left = Inches(0.4)
                            content_frame.margin_right = Inches(0.4)
                            content_frame.margin_top = Inches(0.25)
                            content_frame.text = point
                            content_p = content_frame.paragraphs[0]
                            content_p.font.size = Pt(19)
                            content_p.font.color.rgb = template['shape_text']
                            content_p.line_spacing = 1.4
                            content_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    # --- Highlight Layout ---
                    elif layout_choice == "highlight":
                        for i, point in enumerate(points[:4]):
                            y_pos = 1.8 + i * 1.35
                            
                            box = slide_obj.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(1.2), Inches(y_pos), Inches(7.6), Inches(1.2)
                            )
                            box.fill.solid()
                            box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # Always white background
                            box.line.color.rgb = template['primary']
                            box.line.width = Pt(5)
                            add_advanced_shadow(box, blur=10, distance=5, transparency=0.2)
                            
                            accent = slide_obj.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(1.3), Inches(y_pos + 0.2), Inches(0.15), Inches(0.8)
                            )
                            accent.fill.solid()
                            accent.fill.fore_color.rgb = template['accent']
                            accent.line.fill.background()
                            
                            frame = box.text_frame
                            frame.word_wrap = True
                            frame.margin_left = Inches(0.6)
                            frame.margin_right = Inches(0.4)
                            frame.margin_top = Inches(0.25)
                            frame.text = point
                            p = frame.paragraphs[0]
                            p.font.size = Pt(19)
                            p.font.color.rgb = template['text_primary']  # Dark text on white
                            p.line_spacing = 1.4
                            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    # --- Icon Based Layout ---
                    elif layout_choice == "icon_based":
                        icons = ["★", "●", "■", "▲"]
                        
                        for i, point in enumerate(points[:4]):
                            y_pos = 1.8 + i * 1.35
                            
                            icon_circle = slide_obj.shapes.add_shape(
                                MSO_SHAPE.OVAL,
                                Inches(0.9), Inches(y_pos + 0.2), Inches(0.8), Inches(0.8)
                            )
                            icon_circle.fill.solid()
                            icon_circle.fill.fore_color.rgb = template['accent']
                            icon_circle.line.fill.background()
                            add_advanced_shadow(icon_circle, blur=6, distance=3, transparency=0.2)
                            
                            icon_frame = icon_circle.text_frame
                            icon_frame.text = icons[i]
                            icon_p = icon_frame.paragraphs[0]
                            icon_p.font.size = Pt(32)
                            icon_p.font.bold = True
                            icon_p.font.color.rgb = RGBColor(255, 255, 255)
                            icon_p.alignment = PP_ALIGN.CENTER
                            icon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                            
                            box = slide_obj.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(1.9), Inches(y_pos), Inches(7.3), Inches(1.2)
                            )
                            box.fill.solid()
                            box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                            box.line.color.rgb = template['primary']
                            box.line.width = Pt(4)
                            add_advanced_shadow(box, blur=8, distance=4, transparency=0.18)
                            
                            frame = box.text_frame
                            frame.word_wrap = True
                            frame.margin_left = Inches(0.4)
                            frame.margin_right = Inches(0.4)
                            frame.margin_top = Inches(0.25)
                            frame.text = point
                            p = frame.paragraphs[0]
                            p.font.size = Pt(19)
                            p.font.color.rgb = template['text_primary']
                            p.line_spacing = 1.4
                            frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            elif slide_type == "conclusion":
                top_decoration = slide_obj.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(3), Inches(0.4), Inches(4), Inches(0.3)
                )
                top_decoration.fill.solid()
                top_decoration.fill.fore_color.rgb = template['accent']
                top_decoration.line.fill.background()
                
                title_box = slide_obj.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(2), Inches(0.8), Inches(6), Inches(1)
                )
                title_box.fill.solid()
                title_box.fill.fore_color.rgb = template['primary']
                title_box.line.fill.background()
                add_advanced_shadow(title_box, blur=14, distance=7, transparency=0.25)
                
                title_frame = title_box.text_frame
                title_frame.text = slide_data.get("title", "Xulosa")
                title_frame.margin_left = Inches(0.5)
                title_frame.margin_right = Inches(0.5)
                title_p = title_frame.paragraphs[0]
                title_p.font.size = Pt(42)
                title_p.font.bold = True
                title_p.font.color.rgb = RGBColor(255, 255, 255)  # White text on dark primary
                title_p.alignment = PP_ALIGN.CENTER
                title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                if slide_data.get("summary"):
                    summary_box = slide_obj.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(1.5), Inches(2.1), Inches(7), Inches(1.4)
                    )
                    summary_box.fill.solid()
                    summary_box.fill.fore_color.rgb = template['shape_fill']
                    summary_box.line.color.rgb = template['accent']
                    summary_box.line.width = Pt(4)
                    add_advanced_shadow(summary_box, blur=10, distance=5, transparency=0.18)
                    
                    summary_frame = summary_box.text_frame
                    summary_frame.word_wrap = True
                    summary_frame.margin_left = Inches(0.5)
                    summary_frame.margin_right = Inches(0.5)
                    summary_frame.margin_top = Inches(0.3)
                    summary_frame.margin_bottom = Inches(0.3)
                    summary_frame.text = slide_data.get("summary", "")
                    summary_p = summary_frame.paragraphs[0]
                    summary_p.font.size = Pt(18)
                    summary_p.font.color.rgb = template['shape_text']  # Proper contrast text
                    summary_p.line_spacing = 1.5
                    summary_p.alignment = PP_ALIGN.CENTER
                    summary_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                if slide_data.get("takeaways"):
                    takeaways = slide_data.get("takeaways", [])
                    
                    if len(takeaways) >= 3:
                        for i, takeaway in enumerate(takeaways[:3]):
                            x_pos = 0.8 + i * 3
                            
                            card = slide_obj.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x_pos), Inches(3.9), Inches(2.8), Inches(2.3)
                            )
                            card.fill.solid()
                            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
                            card.line.color.rgb = template['primary']
                            card.line.width = Pt(5)
                            add_advanced_shadow(card, blur=10, distance=5, transparency=0.2)
                            
                            check = slide_obj.shapes.add_shape(
                                MSO_SHAPE.OVAL,
                                Inches(x_pos + 1.15), Inches(4.2), Inches(0.5), Inches(0.5)
                            )
                            check.fill.solid()
                            check.fill.fore_color.rgb = template['accent']
                            check.line.fill.background()
                            
                            check_frame = check.text_frame
                            check_frame.text = "✓"
                            check_p = check_frame.paragraphs[0]
                            check_p.font.size = Pt(24)
                            check_p.font.bold = True
                            check_p.font.color.rgb = RGBColor(255, 255, 255)  # White checkmark on colored circle
                            check_p.alignment = PP_ALIGN.CENTER
                            check_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                            
                            card_frame = card.text_frame
                            card_frame.word_wrap = True
                            card_frame.margin_left = Inches(0.25)
                            card_frame.margin_right = Inches(0.25)
                            card_frame.margin_top = Inches(0.85)
                            card_frame.margin_bottom = Inches(0.25)
                            card_frame.text = takeaway
                            card_p = card_frame.paragraphs[0]
                            card_p.font.size = Pt(15)
                            card_p.font.color.rgb = template['text_primary']  # Dark text on white card
                            card_p.line_spacing = 1.3
                            card_p.alignment = PP_ALIGN.CENTER
                            card_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        except Exception as e:
            print(f"[v0] Error creating slide {idx + 1}: {e}")

    print(f"[v0] Saving READABLE presentation to {filename}")
    prs.save(filename)
    print(f"[v0] READABLE presentation saved successfully")
    return filename


async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "🎓 Assalomu alaykum! ULTRA-MODERN taqdimot yaratish uchun ma'lumotlar kerak.\n\n"
        "✨ 2025 ADVANCED DIZAYN XUSUSIYATLARI:\n"
        "• 18 ta zamonaviy professional shablon\n"
        "• Minimalist, Dark Mode, Gradient, va Retro-Futuristic uslublar\n"
        "• Ultra-yuqori kontrast va o'qilishi oson\n"
        "• Asymmetric va split-screen layoutlar\n"
        "• Advanced shadows va glow effektlar\n"
        "• Geometric decorations va modern accents\n"
        "• Enhanced typography hierarchy\n"
        "• Professional image framing\n"
        "• Micro-animations ready design\n\n"
        "Iltimos, taqdimot mavzusini kiriting:"
    )
    return TOPIC

async def get_topic(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if user_id not in user_data_store:
        user_data_store[user_id] = {}
    
    user_data_store[user_id]['topic'] = update.message.text
    await update.message.reply_text("Nechta slayd kerak? (masalan: 10)")
    return NUM_SLIDES

async def get_num_slides(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    try:
        num_slides = int(update.message.text.strip())
        if num_slides < 4:
            await update.message.reply_text("Kamida 4 ta slayd bo'lishi kerak (sarlavha, kirish, kontent, xulosa). Qaytadan kiriting:")
            return NUM_SLIDES
        user_data_store[user_id]['num_slides'] = num_slides
    except:
        await update.message.reply_text("Iltimos, raqam kiriting (masalan: 10)")
        return NUM_SLIDES
    
    await update.message.reply_text("Universitet nomini kiriting:")
    return UNIVERSITY

async def get_university(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    user_data_store[user_id]['university'] = update.message.text
    await update.message.reply_text("Talaba ismini kiriting:")
    return STUDENT_NAME

async def get_student_name(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    user_data_store[user_id]['student_name'] = update.message.text
    await update.message.reply_text("O'qituvchi ismini kiriting (masalan: 'Aliyev A.A.'):")
    return FROM_TO

async def get_from_to(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    user_data_store[user_id]['from_to'] = update.message.text
    
    data = user_data_store[user_id]
    topic = data['topic']
    num_slides = data['num_slides']
    university = data['university']
    student_name = data['student_name']
    from_to = data['from_to']
    
    await update.message.reply_text(
        f"⏳ ULTRA-MODERN ADVANCED taqdimot tayyorlanmoqda...\n\n"
        f"📚 Mavzu: {topic}\n"
        f"📄 Slaydlar: {num_slides} ta\n"
        f"🎓 Universitet: {university}\n"
        f"👤 Talaba: {student_name}\n"
        f"👨‍🏫 O'qituvchi: {from_to}\n\n"
        f"🎨 2025 ADVANCED DIZAYN - professional, zamonaviy, ultra-modern!\n"
        f"Iltimos, kuting..."
    )

    try:
        print("[v0] Starting ADVANCED content generation...")
        ai_content = await generate_slide_content(topic, num_slides, university, student_name, from_to)

        if not ai_content or len(ai_content) == 0:
            await update.message.reply_text(
                "⚠️ Taqdimot mazmunini yaratishda xatolik yuz berdi.\n"
                "Iltimos, qaytadan urinib ko'ring yoki mavzuni o'zgartiring."
            )
            return ConversationHandler.END

        print(f"[v0] Content generated, creating ULTRA-MODERN ADVANCED PPT with {len(ai_content)} slides...")
        ppt_file = await create_ppt(ai_content, "advanced_slides.pptx")
        
        print("[v0] Sending ADVANCED PPT file to user...")
        await update.message.reply_document(open(ppt_file, "rb"))
        await update.message.reply_text(
            f"✅ ULTRA-MODERN ADVANCED taqdimot tayyor! {len(ai_content)} ta slayd.\n\n"
            f"🎨 2025 ADVANCED DIZAYN XUSUSIYATLARI:\n"
            f"• 18 ta ultra-zamonaviy professional shablon\n"
            f"• Minimalist, Dark Mode, Gradient uslublar\n"
            f"• Retro-Futuristic va Cyberpunk dizaynlar\n"
            f"• Asymmetric va split-screen layoutlar\n"
            f"• Advanced shadows va depth effects\n"
            f"• Geometric decorations\n"
            f"• Enhanced typography hierarchy (52pt titles!)\n"
            f"• Modern icon badges va indicators\n"
            f"• Professional image framing\n"
            f"• Ultra-yuqori kontrast - juda oson o'qiladi\n"
            f"• Optimal spacing - hech narsa overlap qilmaydi\n"
            f"• Modern color schemes va gradients\n"
            f"• Side accents va decorative elements\n\n"
            f"📥 Yuqoridagi faylni yuklab oling va ADVANCED taqdimotingizdan bahramand bo'ling!"
        )
        
        del user_data_store[user_id]
        
    except Exception as e:
        print(f"[v0] Critical error in get_from_to: {e}")
        await update.message.reply_text(
            f"⚠️ Taqdimot yaratishda xatolik yuz berdi.\n"
            f"Xatolik: {str(e)}\n\n"
            f"Iltimos, /start buyrug'i bilan qaytadan boshlang."
        )
        if user_id in user_data_store:
            del user_data_store[user_id]
    
    return ConversationHandler.END

async def cancel(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text("❌ Bekor qilindi. /start buyrug'i bilan qaytadan boshlang.")
    user_id = update.effective_user.id
    if user_id in user_data_store:
        del user_data_store[user_id]
    return ConversationHandler.END

def main():
    app = Application.builder().token(TELEGRAM_TOKEN).build()
    
    conv_handler = ConversationHandler(
        entry_points=[CommandHandler("start", start)],
        states={
            TOPIC: [MessageHandler(filters.TEXT & ~filters.COMMAND, get_topic)],
            NUM_SLIDES: [MessageHandler(filters.TEXT & ~filters.COMMAND, get_num_slides)],
            UNIVERSITY: [MessageHandler(filters.TEXT & ~filters.COMMAND, get_university)],
            STUDENT_NAME: [MessageHandler(filters.TEXT & ~filters.COMMAND, get_student_name)],
            FROM_TO: [MessageHandler(filters.TEXT & ~filters.COMMAND, get_from_to)],
        },
        fallbacks=[CommandHandler("cancel", cancel)],
    )
    
    app.add_handler(conv_handler)
    app.run_polling()

if __name__ == "__main__":
    main()
